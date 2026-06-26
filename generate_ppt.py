import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Set widescreen 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define premium color scheme (AI Recruiter Dashboard dark theme)
    bg_color = RGBColor(12, 12, 20)          # Very deep dark blue-grey
    card_bg_color = RGBColor(22, 22, 38)     # Muted purple-grey card background
    card_border_color = RGBColor(35, 35, 60) # Subtle slate-grey borders
    title_color = RGBColor(0, 223, 216)      # Electric Cyan/Teal
    body_color = RGBColor(255, 255, 255)     # Warm White
    muted_color = RGBColor(160, 160, 176)    # Cool Muted Grey
    accent_color = RGBColor(121, 40, 202)    # Electric Purple
    
    blank_layout = prs.slide_layouts[6]
    
    slides_data = [
        {
            "is_title_slide": True,
            "title": "AI RECRUITER",
            "subtitle": "Intelligent Talent Ingestion & Candidate Alignment System\n\nAutomated Resume Ingestion, Semantic Matching, and Proctor-Secure AI Interviews",
            "meta": "Presenter: Sania\nChandigarh Group of Colleges, Landran\nSubmitted for India Runs by Redrob AI Hackathon (Hack2Skill)"
        },
        {
            "title": "The Core Problem in Recruitment",
            "layout": "3_cols",
            "points": [
                ("Keyword Matching is Broken", "Traditional recruitment software filters profiles using literal keyword matches, missing highly qualified candidates who describe their skills differently."),
                ("Inconsistent Evaluations", "Manual resume reviews are slow and subject to recruiter bias, leading to inconsistent evaluations and missed matches."),
                ("Scale Bottleneck", "Talent acquisition teams spend hours conducting repetitive phone screens and filtering huge stacks of candidates.")
            ]
        },
        {
            "title": "Our Solution Overview",
            "layout": "3_cols",
            "points": [
                ("Double-Layer Screening", "• Layer 1 (Semantic Filtering):\nMatches candidates instantly based on context and meaning, not just keywords.\n\n• Layer 2 (Deep LLM Scoring):\nUses Gemini 2.0 Flash to evaluate fit, list pros/cons, and generate feedback."),
                ("AI Screening Hub", "An interactive, voice-driven screening interview with built-in behavioral telemetry to flag focus-loss or tab-switching in real-time."),
                ("Developer & Stack Flexibility", "Built primarily with a high-performance Python FastAPI backend, supported by an alternative Node.js Express/Supabase stack.")
            ]
        },
        {
            "title": "System Architecture & Workflow",
            "layout": "2x2_grid",
            "points": [
                ("1. Role Creation", "Recruiter creates a job description, which the Google Gemini LLM parses into structured Must-Have and Nice-to-Have skills, experience level, and role summary."),
                ("2. Ingestion & Extraction", "Ingests candidate profiles (via CSV/PDF). Cleans and extracts names, emails, phones, and parsed text automatically."),
                ("3. Vector Matching", "Calculates dense 384-dimensional semantic embeddings. Computes cosine similarity of candidate profiles against the job description requirements."),
                ("4. Dynamic Weighting", "Recruiter can dynamically adjust sliders (Experience vs. Skills vs. Leadership) to instantly recalculate and rank candidate fit scores.")
            ]
        },
        {
            "title": "Layer 1: Semantic Search & Weighting",
            "layout": "3_cols",
            "points": [
                ("Dense Embeddings Model", "Uses local SentenceTransformers ('all-MiniLM-L6-v2') to compute 384-dimensional dense vectors locally to secure high performance."),
                ("Zero-Crash Local Fallback", "If remote servers or libraries are missing, the system uses local NumPy vector calculations and an in-memory cache to ensure zero-downtime evaluation."),
                ("Dynamic Customization", "Features dynamic sliders that allow recruiters to customize candidate rank weighting on-the-fly, bringing total control back to recruiters.")
            ]
        },
        {
            "title": "Layer 2: Deep LLM Evaluation",
            "layout": "3_cols",
            "points": [
                ("Intelligent Fit Scoring", "Gemini 2.0 Flash performs deep text analysis of candidates' career milestones, project details, and technical expertise to assign realistic fit scores."),
                ("Constructive Feedback", "Automatically drafts personalized emails, including structured outreach invites for top matches and customized feedback with learning resources for rejected profiles."),
                ("Outbox Validation Queue", "Queues generated emails in a recruiter outbox for verification, editing, and approval before final transmission via SMTP.")
            ]
        },
        {
            "title": "AI Screening Room & Proctoring",
            "layout": "3_cols",
            "points": [
                ("Biometrics Simulation", "Local camera feed is rendered inside the chat portal with a live face landmark canvas to simulate focus tracking."),
                ("Speech Recognition", "Integrated with the Web Speech API (webkitSpeechRecognition) configured for continuous listening and automatic restart, facilitating natural voice responses."),
                ("Anti-Cheating Telemetry", "Monitors candidate focus, tab-switching, and window blur. Logs warnings immediately to the backend warning API '/api/interviews/proctor/warning/{cand_id}/{job_id}'.")
            ]
        },
        {
            "title": "Workspace Isolation & Compliance Audit",
            "layout": "3_cols",
            "points": [
                ("Secure JWT Authentication", "Access is protected using JSON Web Tokens (JWT) encrypted with bcrypt hashing for secure session logins."),
                ("Recruiter Isolation", "Multi-tenant workspace isolation ensures recruiters only see and manage their own jobs, candidate uploads, and logs."),
                ("Administrative Compliance", "All critical actions (proctoring warnings, uploads, status updates) write directly to a secure compliance Audit Log database table.")
            ]
        },
        {
            "title": "Diagnostic Integrity & Impact",
            "layout": "3_cols",
            "points": [
                ("100% Passed Tests", "All backend diagnostic suites passed: JWT security cryptography, CSV profile ingestion, regex extractors, and workspace isolation."),
                ("Optimized Cost-Efficiency", "Local embedding generation avoids high external API invocation costs, making the platform highly scalable."),
                ("Immediate Ranked Output", "Recruiter dashboard features a live 'Export CSV' button, producing the ranked candidate list in one click for instant submission.")
            ]
        },
        {
            "title": "Conclusion & Future Roadmap",
            "layout": "2_cols",
            "points": [
                ("Functional Summary", "Delivered a complete, working, multi-role recruitment workflow featuring semantic vector search, LLM analysis, and anti-cheating screening."),
                ("Future Extensions", [
                    "Integrate WhatsApp outreach for faster candidate response times.",
                    "Add voice biometrics to verify candidate identity and prevent proxy test-taking.",
                    "Extend multi-lingual interview support using localized Web Speech models."
                ])
            ]
        }
    ]
    
    for s_data in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        
        # Set dark background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        if s_data.get("is_title_slide"):
            # Title slide layout
            # 1. Left accent bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_color
            bar.line.fill.background()
            
            # 2. Main content container card
            panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.9), Inches(5.1))
            panel.fill.solid()
            panel.fill.fore_color.rgb = card_bg_color
            panel.line.color.rgb = card_border_color
            panel.line.width = Pt(1.5)
            
            # 3. Add text box inside panel
            tx_box = slide.shapes.add_textbox(Inches(1.6), Inches(1.6), Inches(10.1), Inches(4.3))
            tf = tx_box.text_frame
            tf.word_wrap = True
            
            # Title
            p = tf.paragraphs[0]
            p.text = s_data["title"]
            p.font.size = Pt(56)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.font.name = "Segoe UI"
            
            # Subtitle
            p2 = tf.add_paragraph()
            p2.text = s_data["subtitle"]
            p2.font.size = Pt(20)
            p2.font.color.rgb = body_color
            p2.font.name = "Segoe UI"
            p2.space_before = Pt(16)
            
            # Meta info
            p3 = tf.add_paragraph()
            p3.text = s_data["meta"]
            p3.font.size = Pt(13.5)
            p3.font.color.rgb = muted_color
            p3.font.name = "Segoe UI"
            p3.space_before = Pt(36)
            
        else:
            # Standard slide layout
            # Add slide title
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = s_data["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = title_color
            p.font.name = "Segoe UI"
            
            # Add divider line
            line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.733), Pt(1.5))
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = card_border_color
            line_shape.line.fill.background()
            
            # Render layout based on layout type
            layout_type = s_data.get("layout")
            points = s_data.get("points", [])
            
            if layout_type == "3_cols":
                # 3 columns layout
                col_width = Inches(3.64)
                col_gap = Inches(0.4)
                top = Inches(1.8)
                height = Inches(4.8)
                
                for idx, pt in enumerate(points):
                    left = Inches(0.8) + idx * (col_width + col_gap)
                    
                    # Draw card background
                    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_width, height)
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg_color
                    card.line.color.rgb = card_border_color
                    card.line.width = Pt(1)
                    
                    # Add top accent border to card (cyan)
                    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_width, Pt(4))
                    top_bar.fill.solid()
                    top_bar.fill.fore_color.rgb = title_color
                    top_bar.line.fill.background()
                    
                    # Text box inside card
                    tx = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), col_width - Inches(0.4), height - Inches(0.5))
                    tf_card = tx.text_frame
                    tf_card.word_wrap = True
                    
                    pt_title, pt_desc = pt
                    p_title = tf_card.paragraphs[0]
                    p_title.text = pt_title
                    p_title.font.size = Pt(18)
                    p_title.font.bold = True
                    p_title.font.color.rgb = title_color
                    p_title.font.name = "Segoe UI"
                    p_title.space_after = Pt(12)
                    
                    if isinstance(pt_desc, list):
                        for sub_pt in pt_desc:
                            p_sub = tf_card.add_paragraph()
                            p_sub.text = "• " + sub_pt
                            p_sub.font.size = Pt(13.5)
                            p_sub.font.color.rgb = body_color
                            p_sub.font.name = "Segoe UI"
                            p_sub.space_before = Pt(4)
                    else:
                        p_desc = tf_card.add_paragraph()
                        p_desc.text = pt_desc
                        p_desc.font.size = Pt(13.5)
                        p_desc.font.color.rgb = body_color
                        p_desc.font.name = "Segoe UI"
                        p_desc.space_before = Pt(4)
                    
            elif layout_type == "2x2_grid":
                # 2x2 grid layout
                col_width = Inches(5.6)
                col_gap = Inches(0.533)
                row_height = Inches(2.2)
                row_gap = Inches(0.4)
                
                for idx, pt in enumerate(points):
                    row = idx // 2
                    col = idx % 2
                    left = Inches(0.8) + col * (col_width + col_gap)
                    top = Inches(1.8) + row * (row_height + row_gap)
                    
                    # Draw card background
                    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_width, row_height)
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg_color
                    card.line.color.rgb = card_border_color
                    card.line.width = Pt(1)
                    
                    # Left accent bar on card (purple)
                    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(4), row_height)
                    left_bar.fill.solid()
                    left_bar.fill.fore_color.rgb = accent_color
                    left_bar.line.fill.background()
                    
                    # Text box inside card
                    tx = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.2), col_width - Inches(0.4), row_height - Inches(0.4))
                    tf_card = tx.text_frame
                    tf_card.word_wrap = True
                    
                    pt_title, pt_desc = pt
                    p_title = tf_card.paragraphs[0]
                    p_title.text = pt_title
                    p_title.font.size = Pt(18)
                    p_title.font.bold = True
                    p_title.font.color.rgb = title_color
                    p_title.font.name = "Segoe UI"
                    p_title.space_after = Pt(6)
                    
                    if isinstance(pt_desc, list):
                        for sub_pt in pt_desc:
                            p_sub = tf_card.add_paragraph()
                            p_sub.text = "• " + sub_pt
                            p_sub.font.size = Pt(13.5)
                            p_sub.font.color.rgb = body_color
                            p_sub.font.name = "Segoe UI"
                            p_sub.space_before = Pt(4)
                    else:
                        p_desc = tf_card.add_paragraph()
                        p_desc.text = pt_desc
                        p_desc.font.size = Pt(13.5)
                        p_desc.font.color.rgb = body_color
                        p_desc.font.name = "Segoe UI"
                        p_desc.space_before = Pt(4)
                    
            elif layout_type == "2_cols":
                # 2 columns layout
                col_width = Inches(5.6)
                col_gap = Inches(0.533)
                top = Inches(1.8)
                height = Inches(4.8)
                
                for idx, pt in enumerate(points):
                    left = Inches(0.8) + idx * (col_width + col_gap)
                    
                    # Draw card background
                    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_width, height)
                    card.fill.solid()
                    card.fill.fore_color.rgb = card_bg_color
                    card.line.color.rgb = card_border_color
                    card.line.width = Pt(1)
                    
                    # Top accent bar (accent_color)
                    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, col_width, Pt(4))
                    top_bar.fill.solid()
                    top_bar.fill.fore_color.rgb = accent_color
                    top_bar.line.fill.background()
                    
                    # Text box
                    tx = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), col_width - Inches(0.6), height - Inches(0.6))
                    tf_card = tx.text_frame
                    tf_card.word_wrap = True
                    
                    pt_title, pt_list = pt
                    p_title = tf_card.paragraphs[0]
                    p_title.text = pt_title
                    p_title.font.size = Pt(20)
                    p_title.font.bold = True
                    p_title.font.color.rgb = title_color
                    p_title.font.name = "Segoe UI"
                    p_title.space_after = Pt(14)
                    
                    if isinstance(pt_list, list):
                        for sub_pt in pt_list:
                            p_sub = tf_card.add_paragraph()
                            p_sub.text = "• " + sub_pt
                            p_sub.font.size = Pt(14.5)
                            p_sub.font.color.rgb = body_color
                            p_sub.font.name = "Segoe UI"
                            p_sub.space_before = Pt(8)
                    else:
                        p_desc = tf_card.add_paragraph()
                        p_desc.text = pt_list
                        p_desc.font.size = Pt(14.5)
                        p_desc.font.color.rgb = body_color
                        p_desc.font.name = "Segoe UI"
                        p_desc.space_before = Pt(6)
                        
    prs.save("AI_Recruiter_Presentation.pptx")
    print("AI_Recruiter_Presentation.pptx generated successfully!")

if __name__ == "__main__":
    create_presentation()
